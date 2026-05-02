"""Generate the Plynos website creation proposal — PPTX (client + template) and PDF.

Run: python3 build_proposal.py
"""

from __future__ import annotations

import os
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------------------
# Brand tokens (Plynos)
# --------------------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x12, 0x20)
BLUE = RGBColor(0x0B, 0x5F, 0xFF)
SLATE = RGBColor(0x5B, 0x64, 0x72)
SOFT_BLUE = RGBColor(0xEA, 0xF2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0B, 0x12, 0x20)
HAIRLINE = RGBColor(0xE5, 0xE7, 0xEB)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
SUBTLE = RGBColor(0xF5, 0xF7, 0xFA)

FONT_HEAD = "Helvetica Neue"
FONT_BODY = "Helvetica Neue"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# --------------------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------------------

def new_pres() -> Presentation:
    pres = Presentation()
    pres.slide_width = Inches(SLIDE_W_IN)
    pres.slide_height = Inches(SLIDE_H_IN)
    return pres


def blank_slide(pres: Presentation, bg: RGBColor = WHITE):
    layout = pres.slide_layouts[6]  # blank
    slide = pres.slides.add_slide(layout)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, pres.slide_width, pres.slide_height
    )
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.shadow.inherit = False
    return slide


def add_text(
    slide,
    text: str,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    font: str = FONT_BODY,
    size: int = 14,
    color: RGBColor = INK,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    spacing: float | None = None,
    char_spacing: float | None = None,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    if char_spacing is not None:
        # python-pptx exposes char spacing via XML
        from pptx.oxml.ns import qn

        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(char_spacing * 100)))
    return tb


def add_paragraphs(
    slide,
    lines,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    font: str = FONT_BODY,
    size: int = 14,
    color: RGBColor = INK,
    bold: bool = False,
    spacing: float = 1.25,
    para_after: int = 6,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(para_after)
        run = p.add_run()
        if isinstance(line, dict):
            run.text = line.get("text", "")
            run.font.name = line.get("font", font)
            run.font.size = Pt(line.get("size", size))
            run.font.color.rgb = line.get("color", color)
            run.font.bold = line.get("bold", bold)
        else:
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
    return tb


def add_rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor | None = None, line_w: float = 0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def add_hairline(slide, x, y, w, color: RGBColor = HAIRLINE):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Emu(7000))
    line.shadow.inherit = False
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_v_hairline(slide, x, y, h, color: RGBColor = HAIRLINE):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Emu(7000), Inches(h))
    line.shadow.inherit = False
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_pill(slide, label: str, x, y, w, h, *, fill=SOFT_BLUE, color=BLUE, size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT_HEAD
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return shape


# --------------------------------------------------------------------------------------
# Common chrome
# --------------------------------------------------------------------------------------

def header_chrome(slide, *, dark: bool = False, page_label: str | None = None):
    """Top wordmark + thin top hairline + page label."""
    color = WHITE if dark else INK
    sub = MUTED if dark else SLATE

    add_text(
        slide,
        "PLYNOS",
        x=0.6, y=0.45, w=2.0, h=0.3,
        font=FONT_HEAD, size=11, color=color, bold=True,
        char_spacing=4,
    )
    if page_label:
        add_text(
            slide,
            page_label,
            x=SLIDE_W_IN - 3.6, y=0.45, w=3.0, h=0.3,
            font=FONT_HEAD, size=10, color=sub, bold=False,
            align=PP_ALIGN.RIGHT, char_spacing=2,
        )


def footer_chrome(slide, *, dark: bool = False):
    color = MUTED if dark else SLATE
    add_hairline(slide, 0.6, SLIDE_H_IN - 0.55, SLIDE_W_IN - 1.2, color=HAIRLINE if not dark else RGBColor(0x1F, 0x29, 0x37))
    add_text(
        slide,
        "Plynos · Website Creation Proposal",
        x=0.6, y=SLIDE_H_IN - 0.42, w=6.0, h=0.3,
        font=FONT_BODY, size=9, color=color,
    )


# --------------------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------------------

@dataclass
class ProposalConfig:
    client_business_name: str
    client_contact_name: str
    client_industry: str
    date: str
    output_pptx: str


def slide_cover(pres, cfg: ProposalConfig):
    slide = blank_slide(pres, bg=NAVY)

    # Subtle blue accent rule on the far left
    add_rect(slide, 0, 0, 0.12, SLIDE_H_IN, fill=BLUE)

    # Wordmark
    add_text(
        slide, "PLYNOS",
        x=0.9, y=0.6, w=3.0, h=0.3,
        font=FONT_HEAD, size=11, color=WHITE, bold=True, char_spacing=4,
    )

    # Eyebrow
    add_text(
        slide, "PROPOSAL",
        x=0.9, y=2.4, w=4.0, h=0.3,
        font=FONT_HEAD, size=11, color=BLUE, bold=True, char_spacing=4,
    )

    # Headline
    add_text(
        slide, "Website Creation\nProposal",
        x=0.9, y=2.8, w=11.0, h=2.2,
        font=FONT_HEAD, size=64, color=WHITE, bold=True, spacing=1.0,
    )

    # Divider
    add_rect(slide, 0.9, 5.15, 0.6, 0.04, fill=BLUE)

    # Subtitle
    add_text(
        slide,
        f"Prepared for {cfg.client_business_name}",
        x=0.9, y=5.35, w=11.0, h=0.5,
        font=FONT_HEAD, size=20, color=WHITE, bold=False,
    )

    # Footer line + meta
    add_rect(slide, 0.9, SLIDE_H_IN - 1.0, SLIDE_W_IN - 1.8, 0.01, fill=RGBColor(0x1F, 0x29, 0x37))
    add_text(
        slide,
        "Prepared by  Harry Davies · Plynos",
        x=0.9, y=SLIDE_H_IN - 0.85, w=8.0, h=0.3,
        font=FONT_BODY, size=11, color=MUTED,
    )
    add_text(
        slide, cfg.date,
        x=SLIDE_W_IN - 4.9, y=SLIDE_H_IN - 0.85, w=4.0, h=0.3,
        font=FONT_BODY, size=11, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def slide_objective(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="01 / OBJECTIVE")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    # Eyebrow
    add_text(
        slide, "OBJECTIVE",
        x=0.6, y=1.5, w=4.0, h=0.3,
        font=FONT_HEAD, size=11, color=BLUE, bold=True, char_spacing=4,
    )

    # Big headline (left)
    add_text(
        slide,
        "A clean, modern website\nthat works as hard\nas the business does.",
        x=0.6, y=2.0, w=7.0, h=3.0,
        font=FONT_HEAD, size=40, color=INK, bold=True, spacing=1.05,
    )

    # Right column body
    add_paragraphs(
        slide,
        [
            f"The goal is to design and build a website for {cfg.client_business_name} that presents the business clearly, builds trust with customers, and makes it easy for them to get in touch.",
            "It should look good on any device, load quickly, and be simple enough that a first-time visitor immediately understands what you do and how to contact you.",
            "Nothing flashy. Nothing complicated. Just a well-built site that does its job, every day.",
        ],
        x=8.0, y=2.05, w=4.8, h=3.5,
        font=FONT_BODY, size=14, color=SLATE, spacing=1.4, para_after=10,
    )

    footer_chrome(slide)


def slide_why(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="02 / WHY THIS MATTERS")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    add_text(
        slide, "Why this matters",
        x=0.6, y=1.4, w=10.0, h=0.7,
        font=FONT_HEAD, size=34, color=INK, bold=True,
    )
    add_text(
        slide,
        "Three practical reasons a professional site pays for itself.",
        x=0.6, y=2.05, w=10.0, h=0.4,
        font=FONT_BODY, size=15, color=SLATE,
    )

    cards = [
        ("01", "Customers check first",
         "Most people look up a business online before they call, message, or visit. A clear, professional site is often the first impression they get."),
        ("02", "A site builds trust",
         "A simple, well-built website signals that the business is real, established, and serious. It reassures customers before they ever pick up the phone."),
        ("03", "One central home",
         "Services, contact details, photos, and credibility, all in one place that you control, not scattered across social profiles or directory listings."),
    ]

    col_w = 3.95
    gap = 0.2
    start_x = 0.6
    top = 2.9
    for i, (num, title, body) in enumerate(cards):
        x = start_x + i * (col_w + gap)
        add_rect(slide, x, top, col_w, 3.6, fill=WHITE, line=HAIRLINE, line_w=0.75)
        add_text(slide, num, x=x + 0.4, y=top + 0.35, w=2.0, h=0.5,
                 font=FONT_HEAD, size=28, color=BLUE, bold=True)
        add_rect(slide, x + 0.4, top + 1.1, 0.4, 0.03, fill=INK)
        add_text(slide, title, x=x + 0.4, y=top + 1.25, w=col_w - 0.8, h=0.6,
                 font=FONT_HEAD, size=18, color=INK, bold=True)
        add_paragraphs(
            slide, [body],
            x=x + 0.4, y=top + 2.0, w=col_w - 0.8, h=1.5,
            font=FONT_BODY, size=12.5, color=SLATE, spacing=1.4, para_after=0,
        )

    footer_chrome(slide)


def slide_package(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="03 / PACKAGE")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    add_text(
        slide, "What's included",
        x=0.6, y=1.4, w=10.0, h=0.7,
        font=FONT_HEAD, size=34, color=INK, bold=True,
    )
    add_text(
        slide,
        f"The standard package, designed around {cfg.client_business_name}.",
        x=0.6, y=2.05, w=10.0, h=0.4,
        font=FONT_BODY, size=15, color=SLATE,
    )

    items = [
        ("Homepage / landing page", "A focused first page that explains what you do at a glance."),
        ("Services section", "A clear breakdown of what you offer, written in plain language."),
        ("About section", "A short, credible introduction to the business and the people behind it."),
        ("Contact section", "Phone, email, address, and a simple contact form, easy to use."),
        ("Mobile-friendly design", "Built to work and feel right on phones, tablets, and desktops."),
        ("Basic SEO setup", "Sensible page titles, descriptions, and structure so your site can be found."),
        ("Fast-loading structure", "Clean code and optimised images so pages open quickly."),
        ("Domain and hosting guidance", "Help choosing and connecting a domain, and getting it live."),
        ("Pre-launch revisions", "A round of adjustments before launch so you're happy with the result."),
    ]

    # Two-column grid of rows with check + label + sub
    col_w = 5.95
    row_h = 0.78
    gap_x = 0.2
    start_x = 0.6
    start_y = 2.65
    for i, (title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + 0.04)

        # Bottom hairline divider per row
        add_hairline(slide, x, y + row_h, col_w)

        # Tick mark in soft blue square
        add_rect(slide, x, y + 0.13, 0.28, 0.28, fill=SOFT_BLUE)
        add_text(slide, "✓", x=x, y=y + 0.09, w=0.28, h=0.35,
                 font=FONT_HEAD, size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

        # Title and description
        add_text(slide, title,
                 x=x + 0.5, y=y + 0.05, w=col_w - 0.5, h=0.33,
                 font=FONT_HEAD, size=13, color=INK, bold=True)
        add_text(slide, body,
                 x=x + 0.5, y=y + 0.38, w=col_w - 0.5, h=0.38,
                 font=FONT_BODY, size=10.5, color=SLATE)

    footer_chrome(slide)


def slide_approach(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="04 / APPROACH")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    add_text(
        slide, "How the site is designed and built",
        x=0.6, y=1.4, w=11.0, h=0.7,
        font=FONT_HEAD, size=34, color=INK, bold=True,
    )
    add_text(
        slide,
        f"Six principles guide the work, from first sketch to launch day for {cfg.client_industry} clients.",
        x=0.6, y=2.05, w=11.0, h=0.4,
        font=FONT_BODY, size=15, color=SLATE,
    )

    items = [
        ("01", "Simple structure",
         "A clean layout that doesn't fight for attention. Information is easy to find."),
        ("02", "Clear messaging",
         "Plain, confident copy that explains what you do without jargon or filler."),
        ("03", "Mobile-first",
         "Designed for the phone first, then scaled up. That's how most people will see it."),
        ("04", "Contact-focused",
         "Every page is built to make getting in touch the obvious next step."),
        ("05", "Easy to read",
         "Generous spacing, strong hierarchy, and copy your customers will actually understand."),
        ("06", "Built around you",
         "Designed around your real services, your real customers, your real brand."),
    ]

    col_w = 4.0
    row_h = 1.95
    gap = 0.15
    start_x = 0.6
    start_y = 2.85
    for i, (num, title, body) in enumerate(items):
        col = i % 3
        row = i // 3
        x = start_x + col * (col_w + gap)
        y = start_y + row * (row_h + 0.15)

        add_rect(slide, x, y, col_w, row_h, fill=SUBTLE)
        # number badge
        add_text(slide, num, x=x + 0.35, y=y + 0.3, w=1.0, h=0.4,
                 font=FONT_HEAD, size=12, color=BLUE, bold=True)
        add_text(slide, title, x=x + 0.35, y=y + 0.7, w=col_w - 0.7, h=0.5,
                 font=FONT_HEAD, size=17, color=INK, bold=True)
        add_paragraphs(
            slide, [body],
            x=x + 0.35, y=y + 1.15, w=col_w - 0.7, h=0.7,
            font=FONT_BODY, size=11.5, color=SLATE, spacing=1.4, para_after=0,
        )

    footer_chrome(slide)


def slide_investment(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="05 / INVESTMENT")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    add_text(
        slide, "Project investment",
        x=0.6, y=1.4, w=10.0, h=0.7,
        font=FONT_HEAD, size=34, color=INK, bold=True,
    )
    add_text(
        slide,
        "One clear price. Paid in two simple stages.",
        x=0.6, y=2.05, w=10.0, h=0.4,
        font=FONT_BODY, size=15, color=SLATE,
    )

    # Left: large price
    add_rect(slide, 0.6, 2.85, 6.0, 3.7, fill=NAVY)
    add_text(slide, "TOTAL",
             x=0.95, y=3.1, w=4.0, h=0.4,
             font=FONT_HEAD, size=11, color=MUTED, bold=True, char_spacing=4)
    add_text(slide, "€600",
             x=0.95, y=3.5, w=5.0, h=1.6,
             font=FONT_HEAD, size=92, color=WHITE, bold=True, spacing=1.0)
    add_text(slide, "+ VAT if applicable",
             x=0.95, y=5.05, w=5.0, h=0.4,
             font=FONT_BODY, size=14, color=MUTED)
    add_rect(slide, 0.95, 5.5, 0.5, 0.03, fill=BLUE)
    add_text(slide,
             "Includes design, development, setup, and launch support.",
             x=0.95, y=5.65, w=4.8, h=0.5,
             font=FONT_BODY, size=12, color=WHITE)

    # Right: payment table
    add_text(slide, "PAYMENT SCHEDULE",
             x=7.0, y=2.85, w=5.0, h=0.4,
             font=FONT_HEAD, size=11, color=BLUE, bold=True, char_spacing=4)

    rows = [
        ("Deposit to begin", "€300", "Paid before work starts."),
        ("Final payment", "€300", "Paid before the site goes live."),
    ]
    y = 3.4
    for i, (label, amt, sub) in enumerate(rows):
        add_hairline(slide, 7.0, y, 5.7)
        add_text(slide, label,
                 x=7.0, y=y + 0.25, w=3.5, h=0.4,
                 font=FONT_HEAD, size=16, color=INK, bold=True)
        add_text(slide, sub,
                 x=7.0, y=y + 0.7, w=4.0, h=0.4,
                 font=FONT_BODY, size=12, color=SLATE)
        add_text(slide, amt,
                 x=10.7, y=y + 0.3, w=2.0, h=0.5,
                 font=FONT_HEAD, size=20, color=INK, bold=True, align=PP_ALIGN.RIGHT)
        y += 1.35
    add_hairline(slide, 7.0, y, 5.7)

    # Total summary line
    add_text(slide, "Total",
             x=7.0, y=y + 0.3, w=3.0, h=0.4,
             font=FONT_HEAD, size=14, color=INK, bold=True)
    add_text(slide, "€600 + VAT",
             x=10.0, y=y + 0.25, w=2.7, h=0.5,
             font=FONT_HEAD, size=18, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)

    footer_chrome(slide)


def slide_timeline(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="06 / TIMELINE")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    add_text(
        slide, "Estimated timeline",
        x=0.6, y=1.4, w=10.0, h=0.7,
        font=FONT_HEAD, size=34, color=INK, bold=True,
    )
    add_text(
        slide,
        "From kick-off to launch, usually within 5–7 business days.",
        x=0.6, y=2.05, w=10.0, h=0.4,
        font=FONT_BODY, size=15, color=SLATE,
    )

    steps = [
        ("Day 0", "Kick-off", "Content and deposit received."),
        ("Day 2–4", "Initial draft", "First design draft delivered for review."),
        ("Day 4–6", "Revisions", "1–3 business days of adjustments and refinements."),
        ("Day 5–7", "Launch", "Site goes live, handover complete."),
    ]

    # Track line
    track_y = 3.85
    track_left = 1.8
    track_right = SLIDE_W_IN - 1.8
    track_w = track_right - track_left
    add_rect(slide, track_left, track_y, track_w, 0.04, fill=HAIRLINE)

    n = len(steps)
    spacing_x = track_w / (n - 1)
    for i, (when, title, body) in enumerate(steps):
        cx = track_left + i * spacing_x
        # dot
        dot_d = 0.28
        add_rect(slide, cx - dot_d / 2, track_y - dot_d / 2 + 0.02, dot_d, dot_d, fill=BLUE)

        # When pill above
        add_text(slide, when,
                 x=cx - 1.0, y=2.85, w=2.0, h=0.4,
                 font=FONT_HEAD, size=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER, char_spacing=3)
        add_text(slide, title,
                 x=cx - 1.4, y=3.2, w=2.8, h=0.5,
                 font=FONT_HEAD, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)

        # Body below
        add_paragraphs(
            slide, [body],
            x=cx - 1.4, y=4.2, w=2.8, h=1.2,
            font=FONT_BODY, size=12, color=SLATE, spacing=1.4, para_after=0, align=PP_ALIGN.CENTER,
        )

    # Note
    add_rect(slide, 0.6, 5.85, SLIDE_W_IN - 1.2, 0.65, fill=SUBTLE)
    add_text(
        slide,
        "Timing depends on how quickly business information, images, and feedback are provided.",
        x=0.85, y=6.0, w=12.0, h=0.5,
        font=FONT_BODY, size=12, color=SLATE, italic=True,
    )

    footer_chrome(slide)


def slide_what_i_need(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="07 / WHAT I NEED")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    add_text(
        slide, f"What I need from {cfg.client_business_name}",
        x=0.6, y=1.4, w=12.0, h=0.7,
        font=FONT_HEAD, size=32, color=INK, bold=True,
    )
    add_text(
        slide,
        "A short list. The faster these arrive, the faster the site goes live.",
        x=0.6, y=2.05, w=11.0, h=0.4,
        font=FONT_BODY, size=15, color=SLATE,
    )

    items = [
        ("Business name", "The official name as you'd like it shown on the site."),
        ("Logo", "If you have one. We can work without it if not."),
        ("Main services", "What you offer, in your own words."),
        ("Short business description", "A few lines about who you are and what you do."),
        ("Contact details", "Phone, email, address, and any social profiles."),
        ("Preferred colours / style", "Optional. Any references or examples you like."),
        ("Photos / images", "Anything you'd like featured. Optional."),
        ("Domain access", "If you already own a domain, the login or registrar."),
    ]

    col_w = 5.95
    row_h = 0.85
    gap_x = 0.2
    start_x = 0.6
    start_y = 2.85
    for i, (title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * row_h

        add_hairline(slide, x, y + row_h - 0.05, col_w)

        # bullet square
        add_rect(slide, x, y + 0.18, 0.18, 0.18, fill=BLUE)

        add_text(slide, title,
                 x=x + 0.4, y=y + 0.05, w=col_w - 0.4, h=0.35,
                 font=FONT_HEAD, size=13.5, color=INK, bold=True)
        add_text(slide, body,
                 x=x + 0.4, y=y + 0.4, w=col_w - 0.4, h=0.4,
                 font=FONT_BODY, size=11, color=SLATE)

    add_rect(slide, 0.6, 6.45, SLIDE_W_IN - 1.2, 0.5, fill=SOFT_BLUE)
    add_text(
        slide,
        "No written content yet? Simple website copy can be drafted for you as part of the package.",
        x=0.85, y=6.55, w=12.0, h=0.4,
        font=FONT_BODY, size=12, color=BLUE, bold=True,
    )

    footer_chrome(slide)


def slide_scope(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="08 / SCOPE")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    add_text(
        slide, "Not included unless agreed separately",
        x=0.6, y=1.4, w=12.0, h=0.7,
        font=FONT_HEAD, size=32, color=INK, bold=True,
    )
    add_text(
        slide,
        "To keep the package clean and the price clear, the items below sit outside the standard scope. Any of them can be added on request.",
        x=0.6, y=2.05, w=11.5, h=0.7,
        font=FONT_BODY, size=14, color=SLATE,
    )

    items = [
        "Paid advertising",
        "Advanced SEO campaigns",
        "Ongoing edits after launch",
        "Professional photography",
        "Complex booking systems",
        "E-commerce functionality",
        "Custom software development",
        "Monthly hosting / domain fees unless agreed",
        "Legal policies unless legal text is provided",
    ]

    # Two-column list, bordered card
    add_rect(slide, 0.6, 3.0, SLIDE_W_IN - 1.2, 3.4, fill=WHITE, line=HAIRLINE, line_w=0.75)

    cols = 2
    inner_pad = 0.4
    col_w = (SLIDE_W_IN - 1.2 - inner_pad * 2) / cols
    row_h = 0.55
    for i, label in enumerate(items):
        col = i % cols
        row = i // cols
        x = 0.6 + inner_pad + col * col_w
        y = 3.2 + row * (row_h + 0.05)

        # dash mark
        add_rect(slide, x, y + 0.27, 0.18, 0.04, fill=SLATE)
        add_text(slide, label,
                 x=x + 0.4, y=y + 0.15, w=col_w - 0.5, h=0.4,
                 font=FONT_HEAD, size=13, color=INK, bold=True)

    # Footer note
    add_text(
        slide,
        "Anything in this list can be added later. Prices for additional work are quoted separately.",
        x=0.6, y=6.55, w=12.0, h=0.4,
        font=FONT_BODY, size=12, color=SLATE, italic=True,
    )

    footer_chrome(slide)


def slide_next_steps(pres, cfg: ProposalConfig):
    slide = blank_slide(pres)
    header_chrome(slide, page_label="09 / NEXT STEPS")
    add_hairline(slide, 0.6, 1.0, SLIDE_W_IN - 1.2)

    add_text(
        slide, "Next steps",
        x=0.6, y=1.4, w=10.0, h=0.7,
        font=FONT_HEAD, size=34, color=INK, bold=True,
    )
    add_text(
        slide,
        "From acceptance to a live website, in five short steps.",
        x=0.6, y=2.05, w=10.0, h=0.4,
        font=FONT_BODY, size=15, color=SLATE,
    )

    steps = [
        ("1", "Confirm acceptance", "By email or WhatsApp, whichever is easier."),
        ("2", "Deposit invoice", "A €300 deposit invoice is sent over."),
        ("3", "Send your details", "Business information, images, and any preferences."),
        ("4", "First draft", "Initial design draft within 2–4 business days."),
        ("5", "Revisions and launch", "Final adjustments, then the site goes live."),
    ]

    start_y = 2.95
    for i, (num, title, body) in enumerate(steps):
        y = start_y + i * 0.6

        # number circle
        add_rect(slide, 0.6, y, 0.42, 0.42, fill=NAVY)
        add_text(slide, num,
                 x=0.6, y=y + 0.05, w=0.42, h=0.4,
                 font=FONT_HEAD, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, title,
                 x=1.2, y=y + 0.05, w=4.5, h=0.4,
                 font=FONT_HEAD, size=15, color=INK, bold=True)
        add_text(slide, body,
                 x=5.8, y=y + 0.07, w=7.0, h=0.4,
                 font=FONT_BODY, size=13, color=SLATE)

    # Closing card
    add_rect(slide, 0.6, 6.05, SLIDE_W_IN - 1.2, 0.85, fill=NAVY)
    add_text(
        slide,
        "Once accepted, work can begin as soon as the deposit has been received.",
        x=0.85, y=6.27, w=10.0, h=0.5,
        font=FONT_HEAD, size=15, color=WHITE, bold=True,
    )
    add_pill(slide, "Ready to start", x=10.65, y=6.27, w=1.8, h=0.4, fill=BLUE, color=WHITE, size=10)

    footer_chrome(slide)


def slide_close(pres, cfg: ProposalConfig):
    slide = blank_slide(pres, bg=NAVY)

    add_rect(slide, 0, 0, 0.12, SLIDE_H_IN, fill=BLUE)

    add_text(slide, "PLYNOS",
             x=0.9, y=0.6, w=3.0, h=0.3,
             font=FONT_HEAD, size=11, color=WHITE, bold=True, char_spacing=4)

    add_text(slide, "Thank you.",
             x=0.9, y=2.6, w=11.0, h=1.4,
             font=FONT_HEAD, size=72, color=WHITE, bold=True)

    add_rect(slide, 0.9, 4.2, 0.6, 0.04, fill=BLUE)

    add_paragraphs(
        slide,
        [
            f"If anything in this proposal needs adjusting for {cfg.client_business_name}, just say. The shape of the work can flex around what's actually useful.",
            "Otherwise, a quick yes by email or WhatsApp is enough to begin.",
        ],
        x=0.9, y=4.4, w=7.9, h=2.2,
        font=FONT_BODY, size=14, color=WHITE, spacing=1.5, para_after=12,
    )

    # Right side: contact card
    add_rect(slide, 9.3, 4.4, 3.4, 1.9, fill=RGBColor(0x12, 0x1B, 0x2D), line=RGBColor(0x1F, 0x29, 0x37))
    add_text(slide, "CONTACT",
             x=9.55, y=4.55, w=3.0, h=0.3,
             font=FONT_HEAD, size=10, color=BLUE, bold=True, char_spacing=4)
    add_text(slide, "Harry Davies",
             x=9.55, y=4.85, w=3.0, h=0.4,
             font=FONT_HEAD, size=16, color=WHITE, bold=True)
    add_text(slide, "Plynos",
             x=9.55, y=5.2, w=3.0, h=0.3,
             font=FONT_BODY, size=12, color=MUTED)
    add_text(slide, "plynos.dev",
             x=9.55, y=5.55, w=3.0, h=0.3,
             font=FONT_BODY, size=12, color=WHITE)
    add_text(slide, "harrycroxforddavies@gmail.com",
             x=9.55, y=5.85, w=3.0, h=0.3,
             font=FONT_BODY, size=12, color=WHITE)

    add_text(slide,
             f"Prepared for {cfg.client_business_name} · {cfg.date}",
             x=0.9, y=SLIDE_H_IN - 0.7, w=11.0, h=0.3,
             font=FONT_BODY, size=10, color=MUTED)


# --------------------------------------------------------------------------------------
# Build orchestration
# --------------------------------------------------------------------------------------

def build_pptx(cfg: ProposalConfig):
    pres = new_pres()
    slide_cover(pres, cfg)
    slide_objective(pres, cfg)
    slide_why(pres, cfg)
    slide_package(pres, cfg)
    slide_approach(pres, cfg)
    slide_investment(pres, cfg)
    slide_timeline(pres, cfg)
    slide_what_i_need(pres, cfg)
    slide_scope(pres, cfg)
    slide_next_steps(pres, cfg)
    slide_close(pres, cfg)
    pres.save(cfg.output_pptx)


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))

    # Client-specific copy (placeholders are kept so the user can fill them in)
    client_cfg = ProposalConfig(
        client_business_name="[Client Business Name]",
        client_contact_name="[Client Contact Name]",
        client_industry="[Client Industry]",
        date="[Date]",
        output_pptx=os.path.join(here, "website_creation_offer_[client_name].pptx"),
    )
    build_pptx(client_cfg)

    # Reusable template
    template_cfg = ProposalConfig(
        client_business_name="[Client Business Name]",
        client_contact_name="[Client Contact Name]",
        client_industry="[Client Industry]",
        date="[Date]",
        output_pptx=os.path.join(here, "website_creation_offer_TEMPLATE.pptx"),
    )
    build_pptx(template_cfg)

    print("PPTX written:")
    print(" -", client_cfg.output_pptx)
    print(" -", template_cfg.output_pptx)
